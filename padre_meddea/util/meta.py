"""Functions to audit information about the mission."""

from pathlib import Path

import astropy.units as u
from astropy.table import Table
from astropy.time import Time
from astropy.timeseries import TimeSeries
from sunpy.net import Fido
from sunpy.net import attrs as a

from padre_meddea import LAUNCH_DATE
from padre_meddea.net import DataType


def _safe_time(value):
    """Return an astropy Time value or None when parsing fails."""
    try:
        return Time(value)
    except Exception:
        return None


def get_all_files(data_type=DataType.housekeeping, level_attr=a.Level.raw):
    """Search for all files of a given data type and level."""
    results = Fido.search(
        a.Time(str(LAUNCH_DATE), str(Time.now()))
        & a.Instrument.meddea
        & data_type
        & level_attr
    )
    return results


def summarize_files(results):
    """Summarize files in a query response by count and total size in bytes."""
    file_ts = queryresponse_to_timeseries(results)
    total_files = len(file_ts)
    if "File Size" in file_ts.colnames:
        total_size = int(file_ts["File Size"].to_value(u.byte).sum())
    else:
        total_size = 0
    return total_files, total_size


def queryresponse_to_timeseries(results, provider="padre"):
    """Convert a QueryResponse provider table to a TimeSeries indexed by file time.

    Parameters
    ----------
    results
        A ``QueryResponse``-like object with provider-key access.
    provider : str
        Provider key in the response object. Defaults to ``padre``.

    Returns
    -------
    astropy.timeseries.TimeSeries
        A time-indexed table containing all non-time columns from the query results.
    """
    provider_results = results[provider]

    if len(provider_results) == 0:
        return TimeSeries(time=Time([], format="unix"))

    parsed_rows = [
        (idx, parsed_time)
        for idx, time_value in enumerate(provider_results["Time"])
        if (parsed_time := _safe_time(time_value)) is not None
    ]

    if len(parsed_rows) == 0:
        return TimeSeries(time=Time([], format="unix"))

    valid_idx = [idx for idx, _ in parsed_rows]
    times = [parsed_time for _, parsed_time in parsed_rows]
    data = {
        colname: [provider_results[colname][idx] for idx in valid_idx]
        for colname in provider_results.colnames
        if colname != "Time"
    }
    if "File Size" in data:
        data["File Size"] = [int(size or 0) for size in data["File Size"]] * u.byte
    ts = TimeSeries(time=Time(times), data=data)
    ts.sort("time")
    return ts


def file_weekly_summary(file_ts):
    """Aggregate a file-level TimeSeries into weekly file totals.

    Parameters
    ----------
    file_ts : astropy.timeseries.TimeSeries
        A file-level time series, typically from ``queryresponse_to_timeseries``.

    Returns
    -------
    astropy.timeseries.TimeSeries
        A weekly ``TimeSeries`` with columns ``Week Number``, ``Total Files``,
        and ``Total Size``.
    """
    if len(file_ts) == 0:
        return TimeSeries(
            time=Time([], format="unix"),
            data={
                "Week Number": [],
                "Total Files": [],
                "Total Size": u.Quantity([], unit=u.Mbyte),
            },
        )

    has_size = "File Size" in file_ts.colnames
    size_values = (
        file_ts["File Size"].to_value(u.byte).astype(int)
        if has_size
        else [0] * len(file_ts)
    )

    iso_years = [int(this_time.strftime("%G")) for this_time in file_ts.time]
    week_numbers = [int(this_time.strftime("%V")) for this_time in file_ts.time]
    iso_weekdays = [int(this_time.strftime("%u")) for this_time in file_ts.time]
    day_starts = [
        Time(f"{this_time.strftime('%Y-%m-%d')}T00:00:00") for this_time in file_ts.time
    ]
    week_starts = [
        day_start - (iso_weekday - 1) * u.day
        for day_start, iso_weekday in zip(day_starts, iso_weekdays, strict=False)
    ]

    ts = TimeSeries(time=file_ts.time, data={"File Size": size_values})
    ts["ISO Year"] = iso_years
    ts["Week Number"] = week_numbers
    ts["Week Start"] = Time(week_starts)
    grouped = Table(ts).group_by(["ISO Year", "Week Number"])

    weekly_times = [group["Week Start"][0] for group in grouped.groups]
    weekly_numbers = [int(group["Week Number"][0]) for group in grouped.groups]
    total_files = [len(group) for group in grouped.groups]
    total_size = [int(sum(group["File Size"])) for group in grouped.groups]

    return TimeSeries(
        time=Time(weekly_times),
        data={
            "Week Number": weekly_numbers,
            "Total Files": total_files,
            "Total Size": u.Quantity(total_size, unit=u.byte).to(u.Mbyte),
        },
    )


def bin_files_by_week(results):
    """Create a weekly summary TimeSeries from query response results."""
    file_ts = queryresponse_to_timeseries(results)
    return file_weekly_summary(file_ts)


def create_summary_powerpoint_slide(
    summary_table,
    output_path,
    title="Weekly File Summary",
    overwrite=True,
):
    """Create a PowerPoint file with one slide containing the summary table.

    Parameters
    ----------
    summary_table : astropy.table.Table or astropy.timeseries.TimeSeries
        The table-like object to render in the slide.
    output_path : str or pathlib.Path
        Destination ``.pptx`` file path.
    title : str
        Slide title text.
    overwrite : bool
        If ``True`` (default), replace an existing file at ``output_path``.
        If ``False``, raise ``FileExistsError`` when the output already exists.

    Returns
    -------
    pathlib.Path
        Path to the written PowerPoint file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        raise ImportError(
            "python-pptx is required to create PowerPoint slides. "
            "Install it with `pip install python-pptx`."
        ) from exc

    table_obj = Table(summary_table)
    output_path = Path(output_path)

    if output_path.exists() and not overwrite:
        raise FileExistsError(f"Output file already exists: {output_path}")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])

    if hasattr(slide.shapes, "title") and slide.shapes.title is not None:
        slide.shapes.title.text = title

    n_rows = len(table_obj) + 1
    n_cols = len(table_obj.colnames)

    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.3)
    height = Inches(5.7)
    ppt_table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    headers = [str(col_name) for col_name in table_obj.colnames]
    rendered_rows = []

    for row in table_obj:
        rendered_row = []
        for col_name in table_obj.colnames:
            value = row[col_name]

            if isinstance(value, Time):
                text_value = value.isot
            elif hasattr(value, "to_value") and hasattr(value, "unit"):
                if col_name == "Total Size":
                    text_value = f"{value.to_value(value.unit):.2f} {value.unit}"
                else:
                    text_value = f"{value.to_value(value.unit):.6g} {value.unit}"
            else:
                text_value = str(value)
            rendered_row.append(text_value)
        rendered_rows.append(rendered_row)

    max_lens = [
        max([len(header)] + [len(row[col_idx]) for row in rendered_rows])
        for col_idx, header in enumerate(headers)
    ]
    total_len = sum(max_lens) if sum(max_lens) > 0 else n_cols

    assigned_width = 0
    for col_idx, col_len in enumerate(max_lens):
        if col_idx < n_cols - 1:
            col_width = int(width * (col_len / total_len))
            assigned_width += col_width
        else:
            col_width = int(width - assigned_width)
        ppt_table.columns[col_idx].width = col_width

    for col_idx, header in enumerate(headers):
        ppt_table.cell(0, col_idx).text = header

    for row_idx, row in enumerate(rendered_rows, start=1):
        for col_idx, text_value in enumerate(row):
            ppt_table.cell(row_idx, col_idx).text = text_value

    presentation.save(output_path)
    return output_path
